import json
import os
from typing import List, Dict, Any
from pathlib import Path
import re
from bs4 import BeautifulSoup
from pptx import Presentation
import nbformat

class ContentExtractor:
    """Extract and chunk content from various file types in instructor_created_data."""
    
    def __init__(self, base_path: str = "instructor_created_data"):
        self.base_path = Path(base_path)
        
    def extract_all_content(self) -> List[Dict[str, Any]]:
        """Extract content from all supported files in the instructor_created_data directory."""
        chunks = []
        
        # Walk through all files in the directory
        for file_path in self.base_path.rglob("*"):
            if file_path.is_file():
                try:
                    file_chunks = self._extract_from_file(file_path)
                    chunks.extend(file_chunks)
                except Exception as e:
                    print(f"⚠️ Error processing {file_path}: {e}")
                    
        print(f"✅ Extracted {len(chunks)} chunks from {len(list(self.base_path.rglob('*')))} files")
        return chunks
    
    def _extract_from_file(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract content from a single file based on its type."""
        suffix = file_path.suffix.lower()
        
        if suffix == '.ipynb':
            return self._extract_jupyter(file_path)
        elif suffix == '.html':
            return self._extract_html(file_path)
        elif suffix == '.json':
            return self._extract_json(file_path)
        elif suffix in ['.ppt', '.pptx']:
            return self._extract_powerpoint(file_path)
        else:
            # Skip unsupported file types
            return []
    
    def _extract_jupyter(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract content from Jupyter notebook by cells."""
        chunks = []
        
        with open(file_path, 'r', encoding='utf-8') as f:
            notebook = nbformat.read(f, as_version=4)
            
        for i, cell in enumerate(notebook.cells):
            if cell.cell_type in ['code', 'markdown'] and cell.source.strip():
                chunk = {
                    'content': cell.source,
                    'source_file': str(file_path.relative_to(self.base_path)),
                    'chunk_type': f'jupyter_cell_{cell.cell_type}',
                    'chunk_id': f'cell_{i}',
                    'metadata': {
                        'cell_type': cell.cell_type,
                        'cell_index': i
                    }
                }
                chunks.append(chunk)
                
        return chunks
    
    def _extract_html(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract content from HTML files by sections."""
        chunks = []
        
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
            
        soup = BeautifulSoup(content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()
        
        # Try to extract by sections (headings)
        sections = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
        
        if sections:
            # Extract content by sections
            for i, heading in enumerate(sections):
                section_content = [heading.get_text()]
                
                # Get content until next heading of same or higher level
                current_level = int(heading.name[1])
                for sibling in heading.next_siblings:
                    if hasattr(sibling, 'name'):
                        if sibling.name and sibling.name.startswith('h'):
                            sibling_level = int(sibling.name[1])
                            if sibling_level <= current_level:
                                break
                        section_content.append(sibling.get_text() if hasattr(sibling, 'get_text') else str(sibling))
                
                content_text = '\n'.join(section_content).strip()
                if content_text:
                    chunk = {
                        'content': content_text,
                        'source_file': str(file_path.relative_to(self.base_path)),
                        'chunk_type': 'html_section',
                        'chunk_id': f'section_{i}',
                        'metadata': {
                            'heading': heading.get_text(),
                            'heading_level': current_level
                        }
                    }
                    chunks.append(chunk)
        else:
            # No clear sections, extract as single chunk
            text = soup.get_text()
            clean_text = re.sub(r'\s+', ' ', text).strip()
            if clean_text:
                chunk = {
                    'content': clean_text,
                    'source_file': str(file_path.relative_to(self.base_path)),
                    'chunk_type': 'html_full',
                    'chunk_id': 'full_content',
                    'metadata': {}
                }
                chunks.append(chunk)
                
        return chunks
    
    def _extract_json(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract content from JSON files."""
        chunks = []
        
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        if isinstance(data, list):
            # Handle list of Q&A items
            for i, item in enumerate(data):
                if isinstance(item, dict):
                    # Combine question and answer for better context
                    content_parts = []
                    if 'student_question' in item:
                        content_parts.append(f"Question: {item['student_question']}")
                    if 'teacher_response' in item:
                        content_parts.append(f"Answer: {item['teacher_response']}")
                    if 'code_snippet' in item:
                        content_parts.append(f"Code: {item['code_snippet']}")
                    
                    content = '\n'.join(content_parts)
                    if content:
                        chunk = {
                            'content': content,
                            'source_file': str(file_path.relative_to(self.base_path)),
                            'chunk_type': 'json_qa',
                            'chunk_id': f'qa_{i}',
                            'metadata': item
                        }
                        chunks.append(chunk)
        elif isinstance(data, dict):
            # Handle single JSON object
            content = json.dumps(data, indent=2)
            chunk = {
                'content': content,
                'source_file': str(file_path.relative_to(self.base_path)),
                'chunk_type': 'json_object',
                'chunk_id': 'full_object',
                'metadata': data
            }
            chunks.append(chunk)
            
        return chunks
    
    def _extract_powerpoint(self, file_path: Path) -> List[Dict[str, Any]]:
        """Extract content from PowerPoint files by slides."""
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            for i, slide in enumerate(prs.slides):
                slide_content = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text.strip())
                
                content = '\n'.join(slide_content)
                if content:
                    chunk = {
                        'content': content,
                        'source_file': str(file_path.relative_to(self.base_path)),
                        'chunk_type': 'powerpoint_slide',
                        'chunk_id': f'slide_{i}',
                        'metadata': {
                            'slide_number': i + 1,
                            'total_slides': len(prs.slides)
                        }
                    }
                    chunks.append(chunk)
                    
        except Exception as e:
            print(f"⚠️ Error reading PowerPoint file {file_path}: {e}")
            
        return chunks

# Example usage and testing
if __name__ == "__main__":
    extractor = ContentExtractor()
    chunks = extractor.extract_all_content()
    
    print(f"\nExtracted {len(chunks)} chunks:")
    for chunk in chunks[:3]:  # Show first 3 chunks
        print(f"\nFile: {chunk['source_file']}")
        print(f"Type: {chunk['chunk_type']}")
        print(f"Content preview: {chunk['content'][:200]}...")